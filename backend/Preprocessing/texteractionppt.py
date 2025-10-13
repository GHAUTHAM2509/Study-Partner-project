# This script extracts text from all slides in a PPTX file and saves it as a .txt file.

from pptx import Presentation
import os,  pathlib, sys

def extract_text_from_pptx(file_path):
    """
    Extracts text from all slides in the given PPTX file.
    - For each slide, collects all text from shapes.
    - Marks the end of each slide.
    - Saves the extracted text as <pptx>.txt.
    """
    prs = Presentation(file_path)
    text_runs = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text_runs.append("\n".join(slide_text))
        text_runs.append(f"slide{i+1} complete")
    text = "\n".join(text_runs)
    pathlib.Path(file_path + ".txt").write_bytes(text.encode())

# Example usage: extract text from a sample PPTX file.
if __name__ == "__main__":
    pptx_file = "Data/ppts/FALLSEM2025-26_VL_BCSE306L_00100_TH_2025-07-25_Intelligent-Agents.pptx"  # Replace with your pptx file path
    extract_text_from_pptx(pptx_file)